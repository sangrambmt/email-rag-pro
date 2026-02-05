import os
from docx import Document
from pypdf import PdfReader
from pptx import Presentation

def extract_text(file_path: str) -> str:
    """Extract text from various file formats"""
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == '.txt':
            return extract_txt(file_path)
        elif ext == '.docx':
            return extract_docx(file_path)
        elif ext == '.pdf':
            return extract_pdf(file_path)
        elif ext in ['.ppt', '.pptx']:
            return extract_pptx(file_path)
        else:
            return ""
    except Exception as e:
        print(f"Error extracting {file_path}: {e}")
        return ""

def extract_txt(file_path: str) -> str:
    """Extract text from .txt files"""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def extract_docx(file_path: str) -> str:
    """Extract text from .docx files"""
    doc = Document(file_path)
    text = []
    
    for para in doc.paragraphs:
        text.append(para.text)
    
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text.append(cell.text)
    
    return '\n'.join(text)

def extract_pdf(file_path: str) -> str:
    """Extract text from .pdf files"""
    reader = PdfReader(file_path)
    text = []
    
    for page in reader.pages:
        text.append(page.extract_text())
    
    return '\n'.join(text)

def extract_pptx(file_path: str) -> str:
    """Extract text from .pptx files"""
    prs = Presentation(file_path)
    text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    
    return '\n'.join(text)