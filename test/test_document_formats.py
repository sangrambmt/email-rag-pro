from pathlib import Path
from docx import Document
from pptx import Presentation

EMAIL_DIR = Path("data/emails")
EMAIL_DIR.mkdir(parents=True, exist_ok=True)

print("Creating test documents...\n")

# DOCX
doc = Document()
doc.add_heading("Sales Proposal", 0)
doc.add_paragraph("Client: TechStart\nAmount: $75,000")
doc.save(EMAIL_DIR / "sales.docx")
print("✅ Created sales.docx")

# PPTX
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Q4 Results"
slide.placeholders[1].text = "Revenue: $2.5M\nGrowth: 25%"
prs.save(EMAIL_DIR / "results.pptx")
print("✅ Created results.pptx")

# TXT
(EMAIL_DIR / "feedback.txt").write_text("Rating: 4/5\nWould recommend: Yes")
print("✅ Created feedback.txt")

print("\n✅ All formats created!")